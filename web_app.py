import streamlit as st
from pptx import Presentation
import io
import os
import tempfile
import subprocess

st.set_page_config(page_title="Sunum Otomasyonu", page_icon="📐")
st.title("📐 Expo Art Design Sunum Otomasyonu")
st.write("Fuar bilgilerini girin, şablonunuzu ve stant render'larını yükleyin.")

# 1. KULLANICI METİN GİRİŞ ALANLARI
st.subheader("📝 Fuar Bilgilerini Girin")
col1, col2 = st.columns(2)
with col1:
    firma_adi = st.text_input("Firma Adı (Örn: Asis)", "Asis")
    fuar_adi = st.text_input("Fuar Adı", "Innotrans 2026")
with col2:
    stand_no = st.text_input("Hall/Stand No", "655")
    olcu = st.text_input("Ölçü", "6.0 x 5.0m(30qm)")

st.divider()

# 2. DOSYA YÜKLEME ALANLARI
st.subheader("📂 Dosyaları Yükleyin")
sablon_dosya = st.file_uploader("Etiketli Şablon Dosyasını Yükleyin", type=["pptx"])
yuklenen_resimler = st.file_uploader("Stant Görsellerini Yükleyin", type=["png", "jpg", "jpeg"], accept_multiple_files=True)

# 3. KESİN HİZALAMA SAĞLAYAN METİN DEĞİŞTİRME FONKSİYONU
def metni_temizle_ve_degistir(paragraph):
    if "{" in paragraph.text and "}" in paragraph.text:
        text = paragraph.text
        text = text.replace("{FIRMA}", firma_adi)
        text = text.replace("{FUAR}", fuar_adi)
        text = text.replace("{STAND}", stand_no)
        text = text.replace("{OLCU}", olcu)
        
        if paragraph.runs:
            paragraph.runs[0].text = text
            p_element = paragraph._p
            for r_element in list(p_element.r_lst[1:]):
                p_element.remove(r_element)
        else:
            paragraph.text = text

def sekilleri_tara_ve_degistir(shapes):
    for shape in shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                metni_temizle_ve_degistir(paragraph)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        metni_temizle_ve_degistir(paragraph)
        elif shape.shape_type == 6: 
            try:
                sekilleri_tara_ve_degistir(shape.shapes)
            except:
                pass

# PDF DÖNÜŞTÜRME FONKSİYONU (Arka planda motoru çalıştırır)
def pptx_to_pdf(pptx_path, output_dir):
    try:
        subprocess.run([
            "libreoffice", "--headless", "--convert-to", "pdf",
            pptx_path, "--outdir", output_dir
        ], check=True)
        return True
    except Exception as e:
        st.error(f"PDF'e dönüştürülürken bir hata oluştu: Lütfen packages.txt dosyasının GitHub'da olduğundan emin olun.")
        return False

# 4. SUNUM OLUŞTURMA İŞLEMİ
if st.button("Sunumu Oluştur", type="primary"):
    if sablon_dosya is not None and len(yuklenen_resimler) > 0:
        with st.spinner("Slaytlar ve PDF hazırlanıyor, bu işlem birkaç saniye sürebilir..."):
            prs = Presentation(sablon_dosya)
            eski_slaytlar = list(prs.slides._sldIdLst)
            
            for layout in prs.slide_layouts:
                sekilleri_tara_ve_degistir(layout.shapes)
            
            yuklenen_resimler = sorted(yuklenen_resimler, key=lambda x: x.name)

            # İlk Slayt
            ilk_duzen = prs.slide_layouts[0]
            slayt_1 = prs.slides.add_slide(ilk_duzen)
            sekilleri_tara_ve_degistir(slayt_1.shapes)
            for shape in slayt_1.placeholders:
                if shape.placeholder_format.type == 18:
                    shape.insert_picture(yuklenen_resimler[0])
                    break

            # Orta Slaytlar
            orta_duzen = prs.slide_layouts[1]
            for resim in yuklenen_resimler[1:]:
                slayt = prs.slides.add_slide(orta_duzen)
                sekilleri_tara_ve_degistir(slayt.shapes)
                for shape in slayt.placeholders:
                    if shape.placeholder_format.type == 18:
                        shape.insert_picture(resim)
                        break

            # Kapanış Slaytı
            bitis_duzeni = prs.slide_layouts[2] 
            slayt_son = prs.slides.add_slide(bitis_duzeni)
            sekilleri_tara_ve_degistir(slayt_son.shapes)

            # Boş sayfaları temizle
            for slide_id in eski_slaytlar:
                prs.slides._sldIdLst.remove(slide_id)

            # Sunumu sunucuda geçici fiziksel bir dosyaya kaydet (PDF motoru okuyabilsin diye)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_pptx:
                prs.save(tmp_pptx.name)
                tmp_pptx_path = tmp_pptx.name

            # PowerPoint çıktısını hafızaya al
            with open(tmp_pptx_path, "rb") as f:
                pptx_bytes = f.read()

            st.success("Tebrikler! Sunum milimetrik hizalamayla başarıyla oluşturuldu.")
            
            # Alt alta değil yan yana iki buton oluşturuyoruz
            buton_kolonu1, buton_kolonu2 = st.columns(2)
            
            with buton_kolonu1:
                st.download_button(
                    label="📊 PowerPoint (.pptx) İndir",
                    data=pptx_bytes,
                    file_name=f"{firma_adi}_Stant_Sunumu.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

            # PDF'E DÖNÜŞTÜRME AŞAMASI
            tmp_pdf_dir = tempfile.gettempdir()
            pdf_basari = pptx_to_pdf(tmp_pptx_path, tmp_pdf_dir)
            
            if pdf_basari:
                base_name = os.path.splitext(os.path.basename(tmp_pptx_path))[0]
                pdf_path = os.path.join(tmp_pdf_dir, base_name + ".pdf")
                
                if os.path.exists(pdf_path):
                    with open(pdf_path, "rb") as f:
                        pdf_bytes = f.read()
                    
                    with buton_kolonu2:
                        st.download_button(
                            label="📄 PDF Olarak İndir",
                            data=pdf_bytes,
                            file_name=f"{firma_adi}_Stant_Sunumu.pdf",
                            mime="application/pdf"
                        )
    else:
        st.warning("Lütfen önce şablonu ve resimleri yüklediğinizden emin olun.")